"""
Blackjack AI — Sunum Oluşturucu
PowerPoint (.pptx) formatında 11 slide'lık Türkçe sunum.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

PLOTS_DIR = "plots"
OUTPUT = "sunum.pptx"

# Colors
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x34, 0x98, 0xDB)
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)
ACCENT_GREEN = RGBColor(0x2E, 0xCC, 0x71)
ACCENT_ORANGE = RGBColor(0xF3, 0x9C, 0x12)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xEC, 0xF0, 0xF1)
DARK_TEXT = RGBColor(0x2C, 0x3E, 0x50)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)


def add_slide():
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    # dark background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    return slide


def add_title(slide, text, top=0.3, size=36, color=WHITE):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.LEFT
    return tf


def add_text(slide, text, left=0.5, top=1.5, width=12, height=5.5, size=18, color=LIGHT_GRAY, bold=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.space_after = Pt(6)
    return tf


def add_image(slide, img_path, left, top, width=None, height=None):
    if os.path.exists(img_path):
        kwargs = {"left": Inches(left), "top": Inches(top)}
        if width:
            kwargs["width"] = Inches(width)
        if height:
            kwargs["height"] = Inches(height)
        slide.shapes.add_picture(img_path, **kwargs)


def add_code(slide, code, left=0.5, top=3.0, width=12, height=3.5):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    # dark code background
    fill = txBox.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x0D, 0x11, 0x17)
    for i, line in enumerate(code.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.name = "Courier New"
        p.font.color.rgb = RGBColor(0xA9, 0xDC, 0x76)
        p.space_after = Pt(2)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1: Kapak
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
add_title(slide, "Blackjack Strateji Öğrenme", top=2.0, size=44)
add_text(slide, "Makine Öğrenmesi ve Pekiştirmeli Öğrenme ile", top=3.0, size=28, color=ACCENT_BLUE)
add_text(slide, "MIS 336 – Uygulamalı Yapay Zeka\nDoç. Dr. Emine Uçar\n\nCan Batu\n18 Mayıs 2026", top=4.2, size=20, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2: Problem & Amaç
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
add_title(slide, "Problem ve Amaç")
add_text(slide, """Blackjack: Oyuncunun kararları sonucu etkileyen nadir casino oyunlarından biri.

Her elde iki seçenek: HIT (kart çek) veya STAND (dur)

Matematiksel olarak kanıtlanmış "basic strategy" ile kasa avantajı %0.5'e düşer.

Amaç: Bu optimal stratejiyi 3 farklı AI yaklaşımıyla öğretmek:
  1. Decision Tree — Yorumlanabilir, kural tabanlı
  2. MLP Neural Network — Backpropagation ile öğrenen sinir ağı
  3. Q-Learning — Etiket olmadan, ödül/ceza ile öğrenen ajan""", top=1.4, size=18)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3: Veri Seti
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
add_title(slide, "Veri Seti")
add_text(slide, """Simülasyon: 300.000 el (6 desteli ayakkabı)
  • Random strateji:  ~%29 kazanma
  • Naive strateji:   ~%39 kazanma
  • Basic strateji:   ~%43 kazanma

Eğitim verisi: 280 benzersiz durum × 500 kopya = 140.000 örnek
  • Özellikler: player_total, dealer_showing, usable_ace
  • Etiket: hit (1) veya stand (0)

Basic strategy fonksiyonundan analitik olarak üretildi.""", top=1.4, size=18)
add_image(slide, f"{PLOTS_DIR}/win_rates.png", left=8.5, top=2.5, width=4.5)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4: Decision Tree
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
add_title(slide, "Model 1: Decision Tree (Karar Ağacı)")
add_text(slide, """• max_depth = 5  |  class_weight = "balanced"
• Doğruluk: %99.3
• Avantaj: Tamamen görselleştirilebilir, her karar izlenebilir
• 201 hata: Sınır bölgesi (oyuncu 12–16 vs krupiye 7–10)""", top=1.4, size=17)
add_code(slide, """dt = DecisionTreeClassifier(
    max_depth=5,
    class_weight="balanced",
    random_state=42
)
dt.fit(X_train, y_train)
# Accuracy: 99.3%""", left=0.5, top=3.2, width=5.5, height=2.8)
add_image(slide, f"{PLOTS_DIR}/decision_tree.png", left=6.5, top=3.0, width=6.5)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5: Decision Tree — Derinlik Deneyi
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
add_title(slide, "Karar Ağacı: max_depth Deneyi")
add_text(slide, """Derinlik arttıkça doğruluk artıyor ama yorumlanabilirlik düşüyor.

  depth=3  →  %89.3  (yetersiz — tüm kuralları öğrenemiyor)
  depth=4  →  %98.2  (iyi ama sınır durumları kaçırıyor)
  depth=5  →  %99.3  ✓ Optimal — kompakt ve yüksek doğruluk
  depth=6  →  %100   (mükemmel ama daha karmaşık ağaç)

Seçim: depth=5 — doğruluk/yorumlanabilirlik dengesi en iyi.""", top=1.4, size=17)
add_image(slide, f"{PLOTS_DIR}/depth_experiment.png", left=7.0, top=2.5, width=5.8)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6: MLP Neural Network
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
add_title(slide, "Model 2: MLP Sinir Ağı (Backpropagation)")
add_text(slide, """Mimari: 3 → 64 → 32 → 1 (2.272 parametre)
Aktivasyon: ReLU  |  Optimizer: Adam  |  Loss: Cross-Entropy
Doğruluk: %100  |  19 epoch'ta yakınsama

Neden %100?
  • Problem deterministik (her durum → tek doğru cevap)
  • Yeterli kapasite (2.272 parametre, 280 kural)
  • Non-linear sınır çizebilir (ReLU)""", top=1.4, size=17)
add_code(slide, """mlp = MLPClassifier(
    hidden_layer_sizes=(64, 32),
    activation="relu",
    solver="adam",
    learning_rate_init=0.001,
    max_iter=200
)
mlp.fit(X_scaled, y_train)""", left=0.5, top=4.0, width=5.5, height=2.5)
add_image(slide, f"{PLOTS_DIR}/mlp_loss_curve.png", left=6.5, top=3.8, width=6.3)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7: Backpropagation Açıklama
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
add_title(slide, "Backpropagation Nasıl Çalışır?")
add_text(slide, """1. İleri Geçiş (Forward Pass)
   Giriş → ağırlıklar × giriş + bias → ReLU → çıkış

2. Kayıp Hesaplama
   Cross-Entropy: L = −[y·log(ŷ) + (1−y)·log(1−ŷ)]

3. Geri Yayılım (Backward Pass)
   Zincir kuralı ile her ağırlığın kayba etkisi hesaplanır:
   ∂L/∂W = (tahmin − gerçek) × önceki katman çıktısı

4. Ağırlık Güncelleme
   W ← W − öğrenme_oranı × ∂L/∂W

5. Tekrar (her epoch'ta tüm veri üzerinden)

Sonuç: Kayıp 0.493 → 0.000085 (19 epoch)""", top=1.3, size=17)
add_image(slide, f"{PLOTS_DIR}/mlp_weights.png", left=0.5, top=5.5, width=12.5)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8: Q-Learning
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
add_title(slide, "Model 3: Q-Learning (Pekiştirmeli Öğrenme)")
add_text(slide, """Etiketli veri YOK — sadece oyun sonucu ile öğrenir.

Q-tablosu: (oyuncu_toplam, krupiye_kart, ace) → [Q_dur, Q_çek]

Güncelleme kuralı (Bellman):
  Q(s,a) ← Q(s,a) + α × [ödül − Q(s,a)]

Hiperparametreler:
  • 500.000 episode  |  α = 0.1  |  ε: 1.0 → 0.05
  • Ödüller: +1.5 blackjack, +1 kazanç, 0 beraberlik, −1 kayıp

Sonuç: %86.2 strateji uyumu  |  %41.9 kazanma oranı""", top=1.3, size=17)
add_image(slide, f"{PLOTS_DIR}/ql_convergence.png", left=7.5, top=3.5, width=5.5)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9: Q-Learning Policy
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
add_title(slide, "Q-Learning: Öğrenilen Strateji vs Optimal")
add_text(slide, """Sol: Basic Strategy (matematiksel optimal)
Sağ: Q-Learning'in öğrendiği strateji

Kırmızı = HIT  |  Yeşil = STAND

%86.2 uyum — etiket görmeden, sadece kazanıp kaybederek öğrendi.""", top=1.3, size=17)
add_image(slide, f"{PLOTS_DIR}/ql_policy.png", left=1.0, top=3.2, width=11.0)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10: Karşılaştırma
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
add_title(slide, "Model Karşılaştırması")
add_text(slide, """Model                Yaklaşım              Doğruluk     Kazanma Oranı
──────────────────────────────────────────────────────────────────────────
Decision Tree        Supervised (depth=5)   %99.3         ~%43
MLP Neural Net       Backpropagation        %100          ~%43
Q-Learning           Reinforcement Learning  %86.2*        %41.9

* Basic strategy ile strateji uyumu

Tüm supervised modeller ~%43 tavan değerine ulaşıyor.
Q-Learning etiketsiz öğrenmesine rağmen tavandan sadece %1.1 uzakta.

Neden %43 tavan?
→ Oyuncu krupiyeden önce oynuyor. Bust olursa krupiye ne çekerse çeksin kaybeder.
→ Bu yapısal kasa avantajı: ~%0.5""", top=1.3, size=16)
add_image(slide, f"{PLOTS_DIR}/model_comparison.png", left=8.0, top=4.0, width=5.0)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11: Heatmap
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
add_title(slide, "Strateji Haritası: Gerçek vs AI")
add_text(slide, """Sol: Matematiksel optimal strateji  |  Sağ: AI (Decision Tree) tahmini
Kırmızı = HIT  |  Yeşil = STAND  |  Neredeyse birebir aynı.""", top=1.3, size=17)
add_image(slide, f"{PLOTS_DIR}/heatmap.png", left=0.5, top=2.5, width=12.0)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12: Sonuç
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
add_title(slide, "Sonuç ve Çıkarımlar")
add_text(slide, """1. Feature engineering en kritik adım.
   Karar anı özellikleri kullanmak: %73 → %99.3 doğruluk

2. Yorumlanabilirlik vs doğruluk trade-off'u gerçek.
   Decision Tree görsel ama %99.3 — MLP %100 ama kara kutu.

3. Backpropagation 19 epoch'ta yakınsadı.
   Loss: 0.493 → 0.000085

4. Q-Learning etiketsiz öğrenmeyi başardı.
   Sıfır bilgiyle başlayıp %86.2 strateji uyumuna ulaştı.

5. %43 tavan yapısal kasa avantajıdır, model sınırlaması değil.

Tüm modeller interaktif Streamlit dashboard'da canlı çalışıyor.""", top=1.3, size=18)
add_text(slide, "Teşekkürler — Sorular?", top=6.5, size=28, color=ACCENT_BLUE, bold=True)

# ═══════════════════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════════════════
prs.save(OUTPUT)
print(f"Sunum oluşturuldu: {OUTPUT}")
print(f"  {len(prs.slides)} slide")
